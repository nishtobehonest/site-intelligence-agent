"""
Generate a 2-slide PPTX for the Site Intelligence Agent demo.

Slide 1 — System Architecture: 4-step pipeline + 3 Chroma collections
Slide 2 — Graceful Degradation: the 3 routing paths with demo examples
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── colour palette ────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0F, 0x17, 0x2A)   # near-black navy – slide background
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xE8, 0xEC, 0xF4)   # light text / subtitles
C_PURPLE    = RGBColor(0x7C, 0x3A, 0xED)   # pipeline steps
C_PURPLE_L  = RGBColor(0xED, 0xE9, 0xFE)   # pipeline step fill (light)
C_TEAL      = RGBColor(0x0E, 0x9A, 0xB5)   # Chroma collections
C_TEAL_L    = RGBColor(0xCC, 0xF5, 0xFA)
C_ORANGE    = RGBColor(0xE6, 0x77, 0x00)   # confidence / partial
C_ORANGE_L  = RGBColor(0xFF, 0xF3, 0xBF)
C_GREEN     = RGBColor(0x2B, 0x8A, 0x3E)   # HIGH
C_GREEN_L   = RGBColor(0xD3, 0xF9, 0xD8)
C_RED       = RGBColor(0xC9, 0x2A, 0x2A)   # LOW
C_RED_L     = RGBColor(0xFF, 0xE3, 0xE3)
C_MAUVE     = RGBColor(0x86, 0x2E, 0x9C)   # LLM / conflict
C_MAUVE_L   = RGBColor(0xF3, 0xD9, 0xFA)
C_GRAY_L    = RGBColor(0xF1, 0xF3, 0xF5)
C_GRAY      = RGBColor(0x86, 0x8E, 0x96)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h,
             fill=None, line_color=None, line_width=Pt(1.5),
             radius=True):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, h
    )
    shape.line.width = line_width
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()   # no line
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if radius:
        # rounded corners via XML
        sp = shape.element
        spPr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
        if spPr is not None:
            from lxml import etree
            prstGeom = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
            if prstGeom is not None:
                prstGeom.set('prst', 'roundRect')
                avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                if avLst is None:
                    avLst = etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                for gd in avLst.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}gd'):
                    avLst.remove(gd)
                etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd',
                                 {'name': 'adj', 'fmla': 'val 20000'})
    return shape


def add_text_box(slide, l, t, w, h, text, font_size=Pt(11),
                 bold=False, color=C_WHITE, align=PP_ALIGN.CENTER,
                 word_wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    # clear default paragraph and set text
    tf.text = ""
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox


def add_arrow(slide, x1, y1, x2, y2, color=C_LIGHT, width=Pt(1.5)):
    """Draw a straight connector arrow from (x1,y1) to (x2,y2)."""
    from pptx.enum.shapes import MSO_CONNECTOR
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # add arrowhead
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = connector.line._ln
    tailEnd = ln.find(qn('a:tailEnd'))
    if tailEnd is None:
        tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'none')
    headEnd = ln.find(qn('a:headEnd'))
    if headEnd is None:
        headEnd = etree.SubElement(ln, qn('a:headEnd'))
    headEnd.set('type', 'arrow')
    headEnd.set('w', 'med')
    headEnd.set('len', 'med')
    return connector


def labeled_box(slide, l, t, w, h,
                title, body=None,
                fill=C_PURPLE_L, border=C_PURPLE,
                title_size=Pt(11), body_size=Pt(9.5),
                title_bold=True, title_color=None, body_color=None):
    """Draw a rounded rect with a title line and optional body lines."""
    add_rect(slide, l, t, w, h, fill=fill, line_color=border, line_width=Pt(1.5))
    tc = title_color or border
    bc = body_color or RGBColor(0x33, 0x33, 0x44)
    if body:
        # title in top portion, body below
        add_text_box(slide, l + Inches(0.08), t + Inches(0.06),
                     w - Inches(0.16), Inches(0.28),
                     title, font_size=title_size, bold=title_bold, color=tc)
        add_text_box(slide, l + Inches(0.1), t + Inches(0.32),
                     w - Inches(0.2), h - Inches(0.38),
                     body, font_size=body_size, bold=False, color=bc,
                     align=PP_ALIGN.LEFT)
    else:
        add_text_box(slide, l + Inches(0.08), t + Inches(0.08),
                     w - Inches(0.16), h - Inches(0.16),
                     title, font_size=title_size, bold=title_bold, color=tc)


# ── SLIDE 1: System Architecture ─────────────────────────────────────────────

def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, C_BG)

    # ── title bar ──────────────────────────────────────────────────────────
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.62),
             fill=C_PURPLE, line_color=None, radius=False)
    add_text_box(slide, Inches(0.2), Inches(0.08), Inches(10), Inches(0.44),
                 "Site Intelligence Agent — System Architecture",
                 font_size=Pt(20), bold=True, color=C_WHITE)
    add_text_box(slide, Inches(10.2), Inches(0.12), Inches(2.9), Inches(0.38),
                 "Site Intelligence Agent · April 2026",
                 font_size=Pt(10), bold=False, color=C_LIGHT, align=PP_ALIGN.RIGHT)

    # ── column x-centres (5 columns: input | step1 | collections | step2+3 | step4+outputs) ──
    # We'll do a top-down layout in 4 "rows"

    ROW_Y = [Inches(0.82), Inches(1.68), Inches(2.62), Inches(3.56), Inches(4.50), Inches(5.44)]
    CX = Inches(6.665)   # horizontal centre of slide

    BOX_W = Inches(4.8)
    BOX_H = Inches(0.70)
    BX = CX - BOX_W / 2   # left x of centred box

    ARROW_X = CX

    # ── [0] Technician Query ───────────────────────────────────────────────
    labeled_box(slide, BX, ROW_Y[0], BOX_W, BOX_H,
                "Technician Query (Natural Language)",
                body="assistant.ask(query)   ·   CLI Interface",
                fill=RGBColor(0xD0, 0xEB, 0xFF),
                border=RGBColor(0x19, 0x71, 0xC2),
                title_color=RGBColor(0x18, 0x64, 0xAB),
                body_color=RGBColor(0x1C, 0x7E, 0xD6))

    add_arrow(slide, ARROW_X, ROW_Y[0] + BOX_H, ARROW_X, ROW_Y[1])

    # ── [1] Step 1: Retrieve ───────────────────────────────────────────────
    labeled_box(slide, BX, ROW_Y[1], BOX_W, BOX_H,
                "STEP 1 — retrieve()   [retriever.py]",
                body="Embed query with all-MiniLM-L6-v2  ·  Search all 3 Chroma collections\nReturn top-k=5 per collection  ·  Sort merged results by cosine similarity",
                fill=C_PURPLE_L, border=C_PURPLE,
                title_color=C_PURPLE)

    # fan arrows to collections
    COLL_Y = ROW_Y[2]
    COLL_H = Inches(0.85)
    COLL_W = Inches(3.5)
    GAP = Inches(0.24)
    COLL_TOTAL = 3 * COLL_W + 2 * GAP
    COLL_X0 = CX - COLL_TOTAL / 2

    cx_osha    = COLL_X0 + COLL_W / 2
    cx_manuals = COLL_X0 + COLL_W + GAP + COLL_W / 2
    cx_jh      = COLL_X0 + 2 * (COLL_W + GAP) + COLL_W / 2

    for cx in [cx_osha, cx_manuals, cx_jh]:
        add_arrow(slide, ARROW_X, ROW_Y[1] + BOX_H, cx, COLL_Y, color=C_TEAL, width=Pt(1.2))

    # ── [2] Three Chroma collections ───────────────────────────────────────
    COLL_DEFS = [
        ("Chroma: osha",
         "29 CFR 1910.147  Lockout/Tagout\n29 CFR 1910.303  Electrical Safety\n566 chunks"),
        ("Chroma: manuals",
         "Carrier 48LC 2017 & 2023\nLennox SL280  ·  Trane XR15\n2,744 chunks"),
        ("Chroma: job_history",
         "synthetic_jobs.json\n50 job records  ·  183 chunks"),
    ]
    for i, (ttl, bdy) in enumerate(COLL_DEFS):
        lx = COLL_X0 + i * (COLL_W + GAP)
        labeled_box(slide, lx, COLL_Y, COLL_W, COLL_H,
                    ttl, body=bdy,
                    fill=C_TEAL_L, border=C_TEAL,
                    title_color=RGBColor(0x0B, 0x72, 0x85),
                    body_color=RGBColor(0x0B, 0x72, 0x85),
                    title_size=Pt(10.5))

    # merge arrows back to centre
    for cx in [cx_osha, cx_manuals, cx_jh]:
        add_arrow(slide, cx, COLL_Y + COLL_H, ARROW_X, ROW_Y[3], color=C_TEAL, width=Pt(1.2))

    # ranked results note
    NOTE_W = Inches(5.2)
    NOTE_H = Inches(0.36)
    add_rect(slide, CX - NOTE_W / 2, ROW_Y[3] - Inches(0.05),
             NOTE_W, NOTE_H,
             fill=RGBColor(0xF1, 0xF3, 0xF5),
             line_color=C_GRAY, line_width=Pt(1), radius=True)
    add_text_box(slide, CX - NOTE_W / 2 + Inches(0.1),
                 ROW_Y[3] - Inches(0.02), NOTE_W - Inches(0.2), NOTE_H,
                 "Merged & ranked  ·  cosine_sim = 1 − (L2_dist² / 2)",
                 font_size=Pt(9), color=C_GRAY, align=PP_ALIGN.CENTER)

    add_arrow(slide, ARROW_X, ROW_Y[3] + NOTE_H - Inches(0.04), ARROW_X, ROW_Y[3] + Inches(0.55))

    # ── [3] Step 2: Score Confidence ──────────────────────────────────────
    S2_Y = ROW_Y[3] + Inches(0.55)
    S2_H = Inches(0.80)
    labeled_box(slide, BX, S2_Y, BOX_W, S2_H,
                "STEP 2 — score_confidence()   [confidence.py]",
                body=("Compare top cosine score against thresholds\n"
                      "detect_conflicts(): cross-collection + within-collection version conflicts"),
                fill=C_ORANGE_L, border=C_ORANGE,
                title_color=C_ORANGE)

    # side annotations
    ANN_W = Inches(2.8)
    ANN_H = Inches(0.90)
    ANN_Y = S2_Y - Inches(0.05)

    # left — conflict detection
    labeled_box(slide, BX - ANN_W - Inches(0.15), ANN_Y, ANN_W, ANN_H,
                "Conflict Detection",
                body=("① Cross-collection: top-3 results\n"
                      "   from ≥2 auth. collections, Δ≤0.15\n"
                      "② Within-collection: 2+ source files,\n"
                      "   best scores Δ≤0.15 (guard ≥0.50)"),
                fill=C_MAUVE_L, border=C_MAUVE,
                title_color=C_MAUVE, title_size=Pt(10),
                body_color=RGBColor(0x5F, 0x3D, 0xC4))

    # right — thresholds
    labeled_box(slide, BX + BOX_W + Inches(0.15), ANN_Y, ANN_W, ANN_H,
                "Confidence Thresholds",
                body=("HIGH:     score ≥ 0.75  AND ≥2 hits ≥0.60\n"
                      "PARTIAL:  0.50 ≤ score < 0.75\n"
                      "          OR conflict detected\n"
                      "LOW:      score < 0.50"),
                fill=RGBColor(0xFF, 0xF9, 0xDB), border=C_ORANGE,
                title_color=C_ORANGE, title_size=Pt(10),
                body_color=RGBColor(0xB3, 0x59, 0x00))

    # ── [4] Step 3: LLM Generate ──────────────────────────────────────────
    S3_Y = S2_Y + S2_H + Inches(0.28)
    S3_H = Inches(0.75)
    labeled_box(slide, BX, S3_Y, BOX_W, S3_H,
                "STEP 3 — llm.generate()   [llm.py]   HIGH / PARTIAL only",
                body=("Gemini 2.5-flash  (swappable via LLM_PROVIDER)\n"
                      "Context: top-5 chunks + system prompt  ·  grounding: \"answer only from context, always cite\"\n"
                      "LOW confidence → this step is skipped entirely (no LLM call)"),
                fill=C_MAUVE_L, border=C_MAUVE,
                title_color=C_MAUVE)

    add_arrow(slide, ARROW_X, S2_Y + S2_H, ARROW_X, S3_Y, color=C_GREEN, width=Pt(1.5))
    # label on arrow
    add_text_box(slide, ARROW_X + Inches(0.06), S2_Y + S2_H + Inches(0.02),
                 Inches(1.8), Inches(0.22),
                 "HIGH / PARTIAL", font_size=Pt(8.5), color=C_GREEN, align=PP_ALIGN.LEFT)

    add_arrow(slide, ARROW_X, S3_Y + S3_H, ARROW_X, ROW_Y[5], width=Pt(1.5))

    # ── [5] Step 4: Route ─────────────────────────────────────────────────
    S4_H = Inches(0.62)
    labeled_box(slide, BX, ROW_Y[5], BOX_W, S4_H,
                "STEP 4 — route()   [degradation.py]",
                body="Format final response by confidence level  ·  Attach source citations",
                fill=C_PURPLE_L, border=C_PURPLE,
                title_color=C_PURPLE)

    # LOW bypass dashed note
    add_text_box(slide, BX - Inches(2.7), S2_Y + Inches(0.18),
                 Inches(1.55), Inches(0.38),
                 "LOW path:\nLLM skipped",
                 font_size=Pt(9), color=C_RED, align=PP_ALIGN.CENTER)

    return slide


# ── SLIDE 2: Graceful Degradation ────────────────────────────────────────────

def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    # title bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.70),
             fill=C_ORANGE, line_color=None, radius=False)
    add_text_box(slide, Inches(0.25), Inches(0.10), Inches(9.5), Inches(0.50),
                 "Graceful Degradation — The Core Differentiator",
                 font_size=Pt(22), bold=True, color=C_WHITE)
    add_text_box(slide, Inches(9.8), Inches(0.14), Inches(3.3), Inches(0.40),
                 "Site Intelligence Agent · April 2026",
                 font_size=Pt(10), color=C_LIGHT, align=PP_ALIGN.RIGHT)

    # subtitle
    add_text_box(slide, Inches(0.4), Inches(0.74), Inches(12.5), Inches(0.32),
                 "Every query is routed to one of three paths based on retrieval confidence. "
                 "False positives (wrong confident answers) are more dangerous than over-escalation.",
                 font_size=Pt(11), color=C_LIGHT, align=PP_ALIGN.LEFT)

    # column layout
    COL_W  = Inches(3.9)
    COL_H  = Inches(5.80)
    COL_Y  = Inches(1.12)
    GAP    = Inches(0.27)
    TOTAL  = 3 * COL_W + 2 * GAP
    COL_X0 = (SLIDE_W - TOTAL) / 2

    COLS = [
        {
            "header": "HIGH CONFIDENCE",
            "header_fill": C_GREEN,
            "box_fill": C_GREEN_L,
            "box_border": C_GREEN,
            "text_color": RGBColor(0x1A, 0x5E, 0x2A),
            "sections": [
                ("Threshold",
                 "top_score ≥ 0.75\nAND ≥ 2 results above 0.60"),
                ("Behaviour",
                 "LLM called with full context block\nAnswer returned with source citations\nescalate: false"),
                ("Conflict check",
                 "No conflict detected\nacross sources"),
                ("Demo query",
                 "\"What are the steps for the\nlockout/tagout energy control\nprocedure?\"\ntop_score = 0.93\n→ [OSHA] 29 CFR 1910.147"),
                ("Eval coverage",
                 "Ground truth set: 94.0%\ncorrect routing (47 / 50)"),
            ]
        },
        {
            "header": "PARTIAL CONFIDENCE",
            "header_fill": C_ORANGE,
            "box_fill": C_ORANGE_L,
            "box_border": C_ORANGE,
            "text_color": RGBColor(0xB3, 0x59, 0x00),
            "sections": [
                ("Threshold",
                 "0.50 ≤ top_score < 0.75\nOR conflict detected between sources"),
                ("Behaviour",
                 "LLM called — answer returned\n⚠ Conflict warning surfaced\nBoth conflicting sources shown\nescalate: true"),
                ("Conflict triggers",
                 "① OSHA vs equipment manual\n② Two manual versions\n   (e.g. Carrier 48LC 2017 vs 2023)"),
                ("Demo query",
                 "\"What is the recommended\nrefrigerant charge pressure\nfor a Carrier rooftop unit?\"\ntop_score = 0.52\n→ 2017 vs 2023 conflict flagged"),
                ("Eval coverage",
                 "Contradictions set: 80.0%\nconflict surfaced (12 / 15)"),
            ]
        },
        {
            "header": "LOW CONFIDENCE",
            "header_fill": C_RED,
            "box_fill": C_RED_L,
            "box_border": C_RED,
            "text_color": RGBColor(0xA6, 0x1E, 0x1E),
            "sections": [
                ("Threshold",
                 "top_score < 0.50\nOR zero results returned"),
                ("Behaviour",
                 "LLM is NOT called\nDeterministic escalation message\nShows closest match + score\nescalate: true"),
                ("Why skip the LLM?",
                 "Calling LLM on low-quality context\nis the hallucination scenario\nthis system is designed to prevent"),
                ("Demo query",
                 "\"What are the repair procedures\nfor a Daikin VRV system\nmodel DX300?\"\ntop_score = 0.31\n→ Equipment not in corpus"),
                ("Eval coverage",
                 "Adversarial set: 45.0% correct\n(near-miss equipment is accepted\nlimitation of dense retrieval)"),
            ]
        },
    ]

    for ci, col in enumerate(COLS):
        lx = COL_X0 + ci * (COL_W + GAP)

        # header strip
        add_rect(slide, lx, COL_Y, COL_W, Inches(0.46),
                 fill=col["header_fill"], line_color=None, radius=True)
        add_text_box(slide, lx + Inches(0.1), COL_Y + Inches(0.07),
                     COL_W - Inches(0.2), Inches(0.32),
                     col["header"], font_size=Pt(14), bold=True, color=C_WHITE)

        # body box
        BODY_Y = COL_Y + Inches(0.50)
        BODY_H = COL_H - Inches(0.54)
        add_rect(slide, lx, BODY_Y, COL_W, BODY_H,
                 fill=col["box_fill"], line_color=col["box_border"],
                 line_width=Pt(1.5), radius=True)

        # sections within body
        sy = BODY_Y + Inches(0.14)
        for sec_title, sec_body in col["sections"]:
            # section title pill
            PILL_H = Inches(0.22)
            add_rect(slide, lx + Inches(0.12), sy,
                     COL_W - Inches(0.24), PILL_H,
                     fill=col["box_border"], line_color=None, radius=True)
            add_text_box(slide, lx + Inches(0.15), sy + Inches(0.02),
                         COL_W - Inches(0.3), PILL_H - Inches(0.04),
                         sec_title, font_size=Pt(8.5), bold=True,
                         color=C_WHITE, align=PP_ALIGN.LEFT)
            sy += PILL_H + Inches(0.04)
            # section body
            lines = sec_body.count("\n") + 1
            sec_h = Inches(0.18 * lines + 0.06)
            add_text_box(slide, lx + Inches(0.16), sy,
                         COL_W - Inches(0.32), sec_h,
                         sec_body, font_size=Pt(10),
                         color=col["text_color"], align=PP_ALIGN.LEFT)
            sy += sec_h + Inches(0.10)

    # bottom bar — key principle
    BAR_Y = Inches(6.98)
    add_rect(slide, Inches(0), BAR_Y, SLIDE_W, Inches(0.52),
             fill=RGBColor(0x1C, 0x27, 0x40), line_color=None, radius=False)
    add_text_box(slide, Inches(0.4), BAR_Y + Inches(0.08), Inches(12.5), Inches(0.36),
                 "Pipeline: retriever.py → confidence.py → llm.py (HIGH/PARTIAL only) → degradation.py  ·  "
                 "Overall eval: 80.0% correct routing across 85 test cases",
                 font_size=Pt(10), color=C_LIGHT, align=PP_ALIGN.CENTER)

    return slide


# ── Main ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    prs = new_prs()
    build_slide1(prs)
    build_slide2(prs)
    out = "Site_Intelligence_Agent_Slides.pptx"
    prs.save(out)
    print(f"Saved: {out}")
