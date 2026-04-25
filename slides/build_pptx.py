"""
Build a PowerPoint deck from the current slide content in slides/index.html.

The PPTX is intentionally presentation-friendly rather than a literal HTML export:
- dark theme to match the Reveal.js deck
- simplified layout for PowerPoint editing
- visible source / citation footer on content slides
- image placeholders on slides where visuals would strengthen the story
- final references slide with URLs
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "site-intelligence-agent.pptx"


BG = RGBColor(15, 23, 42)
PANEL = RGBColor(30, 41, 59)
BORDER = RGBColor(51, 65, 85)
TEXT = RGBColor(226, 232, 240)
MUTED = RGBColor(148, 163, 184)
DIM = RGBColor(100, 116, 139)
ACCENT = RGBColor(96, 165, 250)
HIGH = RGBColor(74, 222, 128)
PARTIAL = RGBColor(250, 204, 21)
LOW = RGBColor(248, 113, 113)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, left, top, width, height, text, font_size=20, color=TEXT,
                bold=False, name=None, align=PP_ALIGN.LEFT, margin=0.05):
    shape = slide.shapes.add_textbox(left, top, width, height)
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    font.name = "Aptos"
    return shape


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.55),
                title, font_size=24, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.color.rgb = BORDER
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.35),
                    subtitle, font_size=10, color=MUTED)


def add_panel(slide, left, top, width, height, title=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    if title:
        add_textbox(slide, left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), Inches(0.25),
                    title, font_size=11, color=MUTED, bold=True)
    return shape


def add_bullets(slide, left, top, width, height, bullets, font_size=14, color=TEXT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Aptos"
    return shape


def add_placeholder(slide, left, top, width, height, title, caption):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(11, 18, 32)
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.5)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.18), width - Inches(0.24), Inches(0.35),
                f"[Image Placeholder] {title}", font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.58), width - Inches(0.36), height - Inches(0.78),
                caption, font_size=11, color=MUTED, align=PP_ALIGN.CENTER)


def add_footer(slide, citation_text, slide_no):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.01)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.color.rgb = BORDER
    add_textbox(slide, Inches(0.5), Inches(6.95), Inches(11.5), Inches(0.3),
                f"Sources: {citation_text}", font_size=8, color=DIM)
    add_textbox(slide, Inches(12.2), Inches(6.93), Inches(0.5), Inches(0.25),
                str(slide_no), font_size=9, color=DIM, align=PP_ALIGN.RIGHT)


def add_badge(slide, left, top, text, fill_rgb, text_rgb):
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(1.05), Inches(0.32))
    badge.fill.solid()
    badge.fill.fore_color.rgb = fill_rgb
    badge.line.color.rgb = fill_rgb
    add_textbox(slide, left, top + Inches(0.015), Inches(1.05), Inches(0.22), text,
                font_size=10, color=text_rgb, bold=True, align=PP_ALIGN.CENTER)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = []

    def new_slide():
        slide = prs.slides.add_slide(blank)
        set_bg(slide)
        slides.append(slide)
        return slide

    # Slide 1
    slide = new_slide()
    add_badge(slide, Inches(0.55), Inches(0.38), "Site Intelligence Agent · April 2026", RGBColor(30, 58, 95), ACCENT)
    add_textbox(slide, Inches(0.7), Inches(1.15), Inches(6.2), Inches(1.1),
                "Site Intelligence\nAgent", font_size=28, color=TEXT, bold=True)
    add_textbox(slide, Inches(0.72), Inches(2.25), Inches(7.0), Inches(0.5),
                "A RAG-Based Context Retrieval System for Frontline Field Workers",
                font_size=16, color=MUTED)
    add_panel(slide, Inches(0.7), Inches(3.0), Inches(6.1), Inches(1.1))
    add_textbox(slide, Inches(0.88), Inches(3.25), Inches(5.7), Inches(0.55),
                '"Field technicians make safety-critical decisions with the wrong information.\nThis system fixes the retrieval and trust problem."', font_size=16, color=TEXT)
    add_placeholder(slide, Inches(7.35), Inches(1.2), Inches(5.15), Inches(3.65),
                    "Hero visual",
                    "Suggested visual: technician in the field using a tablet, with manuals and safety documents fading into a single answer view.")
    add_textbox(slide, Inches(0.9), Inches(5.78), Inches(11.2), Inches(0.42),
                "Nishchay Vishwanath",
                font_size=14, color=TEXT, bold=True)

    # Slide 2
    slide = new_slide()
    add_title(slide, "The Business Problem")
    add_panel(slide, Inches(0.55), Inches(1.0), Inches(6.1), Inches(5.4), "The Gap")
    add_bullets(slide, Inches(0.72), Inches(1.35), Inches(5.7), Inches(1.65), [
        "Technicians make high-stakes decisions under time pressure.",
        "The needed information exists, but it is scattered across manuals, OSHA docs, and office systems.",
        "This is mainly a retrieval and trust problem, not a raw data problem.",
    ], font_size=14)
    add_panel(slide, Inches(0.72), Inches(3.05), Inches(5.7), Inches(1.0), "Who It Affects")
    add_bullets(slide, Inches(0.88), Inches(3.35), Inches(5.3), Inches(0.55), [
        "Technicians: incomplete info and higher error risk",
        "Dispatchers: avoidable escalations and interruptions",
        "Companies: liability, safety violations, and wasted time",
    ], font_size=13)
    add_panel(slide, Inches(6.85), Inches(1.0), Inches(5.95), Inches(2.0), "Current State vs. Target")
    add_bullets(slide, Inches(7.05), Inches(1.35), Inches(5.55), Inches(0.95), [
        "Today: technician calls office, dispatcher searches records, then calls back.",
        "With this system: technician asks on-site and gets a cited answer in seconds.",
    ], font_size=14)
    add_textbox(slide, Inches(7.05), Inches(2.25), Inches(5.55), Inches(0.35),
                "Resolution time: 4 to 12 minutes today vs. under 8 seconds with the assistant",
                font_size=12, color=HIGH, bold=True)
    add_placeholder(slide, Inches(6.85), Inches(3.2), Inches(5.95), Inches(3.2),
                    "Before / after process graphic",
                    "Suggested visual: split image showing dispatcher callback workflow on one side and instant assistant response on the other.")
    add_footer(slide, "Project brief assumptions; field workflow framing from project documentation.", 2)

    # Slide 3
    slide = new_slide()
    add_title(slide, "The AI Problem")
    add_panel(slide, Inches(0.55), Inches(1.0), Inches(5.8), Inches(5.4), "Why AI and Why RAG")
    add_bullets(slide, Inches(0.72), Inches(1.35), Inches(5.45), Inches(1.6), [
        "A technician cannot cross-reference OSHA rules, multiple manual versions, and job history fast enough by hand.",
        "A plain chatbot cannot see private job records and may answer without sources or calibrated confidence.",
        "RAG grounds the answer in retrieved documents so the corpus stays the source of truth.",
    ], font_size=14)
    add_panel(slide, Inches(0.72), Inches(3.2), Inches(5.45), Inches(1.3), "Feasible Data Sources")
    add_bullets(slide, Inches(0.88), Inches(3.55), Inches(5.1), Inches(0.8), [
        "OSHA 29 CFR 1910 and the Field Operations Manual",
        "Carrier, Trane, and Lennox manuals from public manufacturer sites",
        "Synthetic job history generated for the prototype",
    ], font_size=13)
    add_placeholder(slide, Inches(6.6), Inches(1.0), Inches(6.2), Inches(5.4),
                    "RAG concept diagram",
                    "Suggested visual: query -> retrieval over OSHA/manuals/job history -> grounded answer with citations.")
    add_footer(slide, "Lewis et al. 2020; Gao et al. 2023; OSHA directives; manufacturer documentation.", 3)

    # Slide 4
    slide = new_slide()
    add_title(slide, 'Evaluation: "Moving Fast Without Breaking Trust"')
    add_panel(slide, Inches(0.55), Inches(1.0), Inches(7.15), Inches(4.9), "Evaluation Sets")
    add_bullets(slide, Inches(0.75), Inches(1.4), Inches(6.7), Inches(2.1), [
        "Ground truth: known-answer questions to test accuracy and coverage.",
        "Adversarial: prompts with no answer in the corpus to test hallucination resistance.",
        "Contradictions: prompts where sources disagree to test PARTIAL routing and conflict surfacing.",
    ], font_size=14)
    add_panel(slide, Inches(0.75), Inches(3.7), Inches(6.7), Inches(1.5), "Target Metrics")
    add_bullets(slide, Inches(0.95), Inches(4.05), Inches(6.25), Inches(0.9), [
        "Hallucination rate below 2%",
        "Coverage above 80%",
        "Escalation rate intentionally non-zero to reflect uncertainty",
        "High source precision on cited answers",
    ], font_size=13)
    add_placeholder(slide, Inches(7.95), Inches(1.0), Inches(4.85), Inches(4.9),
                    "Metrics chart",
                    "Suggested visual: bar or scorecard for coverage, escalation, hallucination rate, and source precision.")
    add_footer(slide, "Evaluation plan; RAG reliability and calibrated confidence framing.", 4)

    # Slide 5
    slide = new_slide()
    add_title(slide, "Consistency Testing, A/B Design, and Preference Optimization")
    add_panel(slide, Inches(0.55), Inches(1.0), Inches(4.05), Inches(5.7), "Consistency Testing")
    add_bullets(slide, Inches(0.72), Inches(1.35), Inches(3.7), Inches(1.5), [
        "Run the same factual queries multiple times in separate sessions.",
        "Measure output variance and trigger retrieval audits when answers drift too much.",
    ], font_size=13)
    add_panel(slide, Inches(4.65), Inches(1.0), Inches(4.0), Inches(5.7), "A/B Test Design")
    add_bullets(slide, Inches(4.82), Inches(1.35), Inches(3.65), Inches(1.8), [
        "Variant A: answer plus citations only.",
        "Variant B: answer plus HIGH / PARTIAL / LOW confidence language.",
        "Hypothesis: visible uncertainty improves appropriate reliance.",
    ], font_size=13)
    add_panel(slide, Inches(8.7), Inches(1.0), Inches(4.1), Inches(3.0), "Preference Optimization")
    add_bullets(slide, Inches(8.88), Inches(1.35), Inches(3.75), Inches(1.4), [
        "Dispatcher overrides become labeled feedback.",
        "Those logs later become preference pairs for DPO fine-tuning.",
    ], font_size=13)
    add_placeholder(slide, Inches(8.7), Inches(4.2), Inches(4.1), Inches(2.5),
                    "Feedback loop visual",
                    "Suggested visual: system answer -> dispatcher override -> labeled pair -> next model version.")
    add_footer(slide, "Human factors framing from project slides; DPO roadmap described in current slides.", 5)

    # Slide 6
    slide = new_slide()
    add_title(slide, "Architecture: The Four-Layer Pipeline")
    box_w = Inches(2.75)
    xs = [Inches(0.55), Inches(3.28), Inches(6.01), Inches(8.74)]
    titles = [
        ("Layer 1", "Retrieval", "Embed query and retrieve top matches across collections"),
        ("Layer 2", "Confidence", "Assign HIGH, PARTIAL, or LOW from scores and conflicts"),
        ("Layer 3", "Degradation", "Choose answer, warning, or escalation behavior"),
        ("Layer 4", "LLM Generation", "Generate only when context is strong enough"),
    ]
    for x, (num, label, sub) in zip(xs, titles):
        add_panel(slide, x, Inches(1.45), box_w, Inches(1.5))
        add_textbox(slide, x + Inches(0.12), Inches(1.63), box_w - Inches(0.24), Inches(0.2), num, font_size=10, color=DIM)
        add_textbox(slide, x + Inches(0.12), Inches(1.9), box_w - Inches(0.24), Inches(0.25), label, font_size=15, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.16), Inches(2.22), box_w - Inches(0.32), Inches(0.45), sub, font_size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_panel(slide, Inches(0.55), Inches(3.35), Inches(8.25), Inches(2.5), "Confidence Routing")
    add_bullets(slide, Inches(0.75), Inches(3.72), Inches(7.8), Inches(1.6), [
        "HIGH: strong match and no conflict -> cited answer",
        "PARTIAL: weaker evidence or source conflict -> answer plus warning and surfaced conflict",
        "LOW: weak or missing context -> skip the LLM and escalate programmatically",
    ], font_size=13)
    add_placeholder(slide, Inches(9.0), Inches(3.35), Inches(3.8), Inches(2.5),
                    "Pipeline diagram",
                    "Suggested visual: horizontal flow with color-coded HIGH / PARTIAL / LOW routes.")
    add_footer(slide, "Project architecture from README, CLAUDE.md, and slide deck content.", 6)

    # Slide 7
    slide = new_slide()
    add_title(slide, "Context and Memory Architecture")
    add_panel(slide, Inches(0.55), Inches(1.0), Inches(6.1), Inches(5.6), "What Goes Into Context")
    add_bullets(slide, Inches(0.72), Inches(1.35), Inches(5.75), Inches(1.8), [
        "System prompt: cite sources, do not fabricate, flag safety-sensitive procedures.",
        "Retrieved chunks: top results with collection, document name, and similarity score.",
        "Technician query: equipment, job context, and natural language question.",
    ], font_size=14)
    add_panel(slide, Inches(0.72), Inches(3.45), Inches(5.75), Inches(1.6), "Memory Types")
    add_bullets(slide, Inches(0.9), Inches(3.82), Inches(5.4), Inches(1.0), [
        "Working memory: retrieved chunks for this single query",
        "Institutional memory: job_history collection",
        "Semantic memory: OSHA and manuals collections",
        "Session memory: stateless in v1, possible in v2",
    ], font_size=13)
    add_placeholder(slide, Inches(6.9), Inches(1.0), Inches(5.9), Inches(5.6),
                    "Memory architecture visual",
                    "Suggested visual: stacked memory model showing query context, semantic docs, and institutional job history.")
    add_footer(slide, "Project system prompt and collection design from repository docs and code.", 7)

    # Slide 8
    slide = new_slide()
    add_title(slide, "Guardrails and AI Concept Map")
    add_panel(slide, Inches(0.55), Inches(1.0), Inches(5.8), Inches(5.6), "Guardrails")
    add_bullets(slide, Inches(0.72), Inches(1.35), Inches(5.45), Inches(2.2), [
        "Citation requirement on every answer",
        "Prompt-level hallucination prevention",
        "Explicit safety flagging for risky procedures",
        "Scope boundary for off-topic or weak-context questions",
    ], font_size=14)
    add_panel(slide, Inches(0.72), Inches(4.2), Inches(5.45), Inches(1.3), "Design Principle")
    add_textbox(slide, Inches(0.9), Inches(4.55), Inches(5.1), Inches(0.55),
                "In a safety-critical context, a confident wrong answer is worse than an unnecessary escalation.",
                font_size=13, color=PARTIAL)
    add_placeholder(slide, Inches(6.55), Inches(1.0), Inches(6.25), Inches(5.6),
                    "Concept map",
                    "Suggested visual: map RAG, memories, tools/MCP, fine-tuning, and preference optimization to where they live in the system.")
    add_footer(slide, "Guardrails and concept mapping from current deck content and project docs.", 8)

    # Slide 9
    slide = new_slide()
    add_title(slide, "Cost Model for 100 Users")
    add_panel(slide, Inches(0.55), Inches(1.0), Inches(5.6), Inches(5.6), "Operating Assumptions")
    add_bullets(slide, Inches(0.72), Inches(1.35), Inches(5.25), Inches(2.0), [
        "100 technicians, 10 queries per day, 250 workdays per year",
        "About 250,000 annual queries, with LOW-routed cases skipping the LLM",
        "Local Chroma and local embeddings keep infrastructure cost down",
    ], font_size=14)
    add_panel(slide, Inches(6.35), Inches(1.0), Inches(6.45), Inches(3.3), "Monthly Cost Snapshot")
    add_bullets(slide, Inches(6.55), Inches(1.35), Inches(6.05), Inches(1.8), [
        "Claude API input and output: roughly $150 per month combined in the current assumptions",
        "Local vector DB and local embeddings: effectively $0",
        "Hosting estimate: roughly $75 to $100 per month",
        "Total estimate: about $280 to $350 per month",
    ], font_size=13)
    add_placeholder(slide, Inches(6.35), Inches(4.55), Inches(6.45), Inches(2.05),
                    "Cost chart",
                    "Suggested visual: stacked bar for API, hosting, and infrastructure cost categories.")
    add_footer(slide, "Cost model; Claude pricing should be verified against current Anthropic pricing page before final delivery.", 9)

    # Slide 10
    slide = new_slide()
    add_title(slide, "Claude vs. Llama 3 70B: Break-Even Analysis")
    add_panel(slide, Inches(0.55), Inches(1.0), Inches(7.25), Inches(5.6), "Comparison")
    add_bullets(slide, Inches(0.72), Inches(1.35), Inches(6.9), Inches(2.5), [
        "Claude wins early on simplicity, no GPU infrastructure, and immediate quality.",
        "Self-hosted Llama becomes more attractive at higher query volume and after collecting enough labeled preference data.",
        "The current architecture stays model-agnostic so migration is operationally simpler.",
    ], font_size=14)
    add_textbox(slide, Inches(0.75), Inches(4.6), Inches(6.9), Inches(0.55),
                "Break-even claim in the current deck: roughly 55,000 to 60,000 queries per month (about 550 users).",
                font_size=12, color=HIGH, bold=True)
    add_placeholder(slide, Inches(8.0), Inches(1.0), Inches(4.8), Inches(5.6),
                    "Break-even chart",
                    "Suggested visual: two cost lines crossing over as usage grows, with the break-even point labeled.")
    add_footer(slide, "Cost model plus self-hosted vs API tradeoff framing from the current deck.", 10)

    # Slide 11
    slide = new_slide()
    add_title(slide, "Revenue Model and Key Tradeoffs")
    add_panel(slide, Inches(0.55), Inches(1.0), Inches(6.0), Inches(5.6), "Revenue Options")
    add_bullets(slide, Inches(0.72), Inches(1.35), Inches(5.65), Inches(1.7), [
        "Per-seat SaaS pricing",
        "Per-query pricing",
        "Enterprise fleet license",
    ], font_size=14)
    add_panel(slide, Inches(0.72), Inches(3.15), Inches(5.65), Inches(2.0), "Tradeoffs")
    add_bullets(slide, Inches(0.9), Inches(3.52), Inches(5.3), Inches(1.35), [
        "Precision vs. escalation rate",
        "Corpus freshness vs. operating cost",
        "API dependency vs. on-prem control",
        "Fine-tuning ROI vs. timeline",
    ], font_size=13)
    add_placeholder(slide, Inches(6.8), Inches(1.0), Inches(6.0), Inches(3.0),
                    "ROI graphic",
                    "Suggested visual: value equation linking fewer incidents, faster resolution, and lower dispatcher load to customer ROI.")
    add_panel(slide, Inches(6.8), Inches(4.3), Inches(6.0), Inches(2.3), "Core Thesis")
    add_textbox(slide, Inches(7.0), Inches(4.75), Inches(5.6), Inches(1.05),
                "The product is valuable because it answers quickly when it should, and escalates clearly when it should not pretend to know.",
                font_size=14, color=TEXT)
    add_footer(slide, "Business model framing from the current deck.", 11)

    # Slide 12
    slide = new_slide()
    add_title(slide, "Sources, Citations, and References")
    add_panel(slide, Inches(0.55), Inches(1.0), Inches(12.25), Inches(5.9), "Reference List")
    refs = [
        "Lewis, P. et al. Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks. NeurIPS 2020. https://arxiv.org/abs/2005.11401",
        "Gao, Y. et al. Retrieval-Augmented Generation for Large Language Models: A Survey. 2023. https://arxiv.org/abs/2312.10997",
        "OSHA Field Operations Manual. U.S. Department of Labor. https://www.osha.gov/enforcement/directives/cpl-02-00-164",
        "OSHA 29 CFR 1910.147 Lockout/Tagout. https://www.osha.gov/laws-regs/regulations/standardnumber/1910/1910.147",
        "OSHA 29 CFR 1910.303 Electrical, General. https://www.osha.gov/laws-regs/regulations/standardnumber/1910/1910.303",
        "Anthropic Claude model and pricing pages. Verify latest pricing before presenting cost claims.",
        "Public manufacturer documentation for Carrier, Trane, and Lennox equipment manuals.",
        "Project repository sources: README.md, CLAUDE.md, src/assistant.py, src/retriever.py, src/confidence.py, src/degradation.py, and slides/index.html",
    ]
    add_bullets(slide, Inches(0.78), Inches(1.35), Inches(11.8), Inches(4.9), refs, font_size=12)
    add_textbox(slide, Inches(0.78), Inches(6.15), Inches(11.8), Inches(0.35),
                "Note: pricing and any vendor-specific claims should be re-verified before final submission.", font_size=10, color=PARTIAL)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_deck()
